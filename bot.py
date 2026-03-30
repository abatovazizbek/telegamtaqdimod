import asyncio
import logging
import io
import re
import requests

from aiogram import Bot, Dispatcher, F, types
from aiogram.filters import Command
from aiogram.utils.keyboard import InlineKeyboardBuilder

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# ===================== YANGI GEMINI SDK =====================
from google import genai

# ===================== SOZLAMALAR =====================
TOKEN = "8128500951:AAFsgE6uq8eX2kY8_yxFnCLajzrEE3p7EtY"
GEMINI_API_KEY = "AIzaSyBzg_66XVCdCX2JYRObFNVOZYAkpHcNptM"

logging.basicConfig(level=logging.INFO)

bot = Bot(token=TOKEN)
dp = Dispatcher()

# ===================== GEMINI CLIENT =====================
client = genai.Client(api_key=GEMINI_API_KEY)

def generate_text(prompt: str) -> str | None:
    """Kvota tejash uchun asosan Flash modelidan foydalanamiz"""
    models = [
        "gemini-2.5-flash",       # Eng tavsiya etiladi (tez + yaxshi kvota)
        "gemini-2.5-flash-lite",  # Eng yuqori kvota free tierda
        "gemini-2.5-pro"          # Oxirgi variant (tezroq tugaydi)
    ]
    
    for model_name in models:
        try:
            logging.info(f"Model sinab ko‘rilmoqda: {model_name}")
            
            response = client.models.generate_content(
                model=model_name,
                contents=prompt
            )
            
            if response and response.text:
                logging.info(f"✅ Muvaffaqiyatli javob: {model_name}")
                return response.text
                
        except Exception as e:
            error_str = str(e).lower()
            logging.error(f"{model_name} xatosi: {e}")
            
            if "429" in error_str or "resource_exhausted" in error_str or "quota" in error_str:
                logging.warning(f"⚠️ {model_name} kvotasi tugagan, keyingi modelga o'tilmoqda...")
                continue
            else:
                # Boshqa xato bo'lsa ham keyingi modelni sinaymiz
                continue
    
    logging.error("Barcha modellar xato berdi yoki kvota tugadi.")
    return None


user_data = {}

# ===================== YORDAMCHI FUNKSIYALAR =====================
def set_slide_background(slide, color_rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color_rgb

def get_image(query: str):
    """Unsplash orqali rasm olish"""
    try:
        safe_query = re.sub(r'[^a-zA-Z0-9]', '-', query.lower())[:50]
        url = f"https://source.unsplash.com/featured/800x600/?{safe_query}"
        if requests.head(url, timeout=5).status_code == 200:
            return url
    except:
        pass
    return None

def create_pptx(raw_text: str, bg_type: str):
    prs = Presentation()
    
    bg_colors = {
        "blue": RGBColor(0, 32, 96),
        "dark": RGBColor(33, 33, 33),
        "white": RGBColor(255, 255, 255)
    }
    text_colors = {
        "blue": RGBColor(255, 255, 255),
        "dark": RGBColor(255, 255, 255),
        "white": RGBColor(0, 0, 0)
    }
    
    sel_bg = bg_colors.get(bg_type, bg_colors["white"])
    sel_txt = text_colors.get(bg_type, text_colors["white"])
    
    sections = re.split(r'###', raw_text)
    
    for section in sections:
        section = section.strip()
        if len(section) < 15:
            continue
            
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_slide_background(slide, sel_bg)
        
        lines = section.split('\n')
        if not lines:
            continue
            
        # Sarlavha
        title = slide.shapes.title
        title.text = lines[0].replace("*", "").strip()
        title.text_frame.paragraphs[0].font.color.rgb = sel_txt
        title.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Matn
        body = slide.placeholders[1]
        for line in lines[1:]:
            clean = line.strip().replace("*", "").replace("- ", "")
            if clean:
                p = body.text_frame.add_paragraph()
                p.text = clean
                p.font.size = Pt(18)
                p.font.color.rgb = sel_txt
        
        # Rasm qo'shish
        img_url = get_image(title.text)
        if img_url:
            try:
                img_data = requests.get(img_url, timeout=8).content
                slide.shapes.add_picture(io.BytesIO(img_data), Inches(5.3), Inches(1.6), width=Inches(4.0))
            except:
                pass
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ===================== BOT HANDLERLARI =====================
@dp.message(Command("start"))
async def cmd_start(message: types.Message):
    await message.answer("✍️ Taqdimot mavzusini yozing\n(masalan: Sun'iy intellekt, O'zbekiston tarixi, Python dasturlash)")

@dp.message(F.text)
async def start_flow(message: types.Message):
    if message.text.startswith('/'):
        return
    user_data[message.from_user.id] = message.text
    
    kb = InlineKeyboardBuilder()
    kb.row(
        types.InlineKeyboardButton(text="🔵 Ko'k", callback_data="bg_blue"),
        types.InlineKeyboardButton(text="⚫️ Qora", callback_data="bg_dark"),
        types.InlineKeyboardButton(text="⚪️ Oq", callback_data="bg_white")
    )
    await message.answer("🎨 Fon rangini tanlang:", reply_markup=kb.as_markup())

@dp.callback_query(F.data.startswith("bg_"))
async def handle_callback(callback: types.CallbackQuery):
    bg_type = callback.data.split("_")[1]
    topic = user_data.get(callback.from_user.id)
    
    if not topic:
        await callback.message.answer("Mavzu topilmadi. /start bilan qaytadan boshlang.")
        return
    
    await callback.message.edit_text(f"⏳ '{topic}' taqdimoti tayyorlanmoqda...\n\nBu biroz vaqt olishi mumkin.")

    try:
        prompt = f"Write 5 clear, educational and well-structured slides about '{topic}' in Uzbek language. Use simple and understandable language for students. Separate each slide with ###."
        
        res_text = generate_text(prompt)
        
        if not res_text:
            raise Exception("Gemini javob bermadi. API kvotasi tugagan bo'lishi mumkin. Keyinroq urinib ko'ring yoki yangi API key oling.")

        pptx_buffer = create_pptx(res_text, bg_type)
        file = types.BufferedInputFile(pptx_buffer.read(), filename=f"{topic[:50]}.pptx")
        
        await callback.message.answer_document(file, caption=f"✅ '{topic}' taqdimoti tayyor bo‘ldi!")
        await callback.message.delete()
        
    except Exception as e:
        logging.error(f"Umumiy xato: {e}")
        await callback.message.answer(f"❌ Xato yuz berdi:\n{str(e)[:500]}\n\nKeyinroq qayta urinib ko‘ring.")

async def main():
    await bot.delete_webhook(drop_pending_updates=True)
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
